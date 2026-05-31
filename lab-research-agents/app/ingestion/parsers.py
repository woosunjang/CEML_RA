"""
Lab Research Agents — File Parsers

Parse PDF, DOCX, PPTX, TXT, and MD files into a list of page/text dicts.
Each parser returns: list[dict] where dict = {"page": int | None, "text": str}
"""

from pathlib import Path
from typing import Union


def parse_pdf(file_path: Union[str, Path]) -> list[dict]:
    """Extract text from a PDF file, page by page, using PyMuPDF."""
    import fitz  # PyMuPDF

    pages: list[dict] = []
    doc = fitz.open(str(file_path))
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()
            if text:
                pages.append({"page": page_num + 1, "text": text})
    finally:
        doc.close()
    return pages


def parse_docx(file_path: Union[str, Path]) -> list[dict]:
    """Extract text from a DOCX file as a single block (page=None)."""
    from docx import Document

    doc = Document(str(file_path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    full_text = "\n\n".join(paragraphs)
    return [{"page": None, "text": full_text}]


def parse_pptx(file_path: Union[str, Path]) -> list[dict]:
    """Extract text from a PPTX file, slide by slide (page=slide_number)."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides: list[dict] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)
        if texts:
            slides.append({"page": slide_num, "text": "\n".join(texts)})
    return slides


def parse_text(file_path: Union[str, Path]) -> list[dict]:
    """Read a plain text or markdown file as a single block (page=None)."""
    path = Path(file_path)
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    if not text:
        return []
    return [{"page": None, "text": text}]


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_PARSER_MAP = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".txt": parse_text,
    ".md": parse_text,
}


def parse_file(file_path: Union[str, Path]) -> list[dict]:
    """
    Dispatch to the appropriate parser based on file extension.

    Returns:
        list[dict] — each dict has keys "page" (int | None) and "text" (str).

    Raises:
        ValueError: if the file extension is not supported.
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    parser = _PARSER_MAP.get(ext)
    if parser is None:
        supported = ", ".join(sorted(_PARSER_MAP.keys()))
        raise ValueError(
            f"Unsupported file type: '{ext}'. Supported types: {supported}"
        )
    return parser(path)
