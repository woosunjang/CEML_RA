"""
Presentation Agent — PPTX Builder

Converts JSON slide structure to actual PowerPoint (.pptx) files.
Supports multiple themes, auto font detection (Noto Sans / Noto Sans KR),
and optional image insertion.
"""

import json
import logging
import re
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Theme system
# ---------------------------------------------------------------------------
THEMES = {
    "dark_academic": {
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "title": RGBColor(0xE8, 0xE8, 0xF0),
        "body": RGBColor(0xD0, 0xD0, 0xE0),
        "accent": RGBColor(0x81, 0x8C, 0xF8),
        "subtitle": RGBColor(0xA0, 0xA0, 0xC0),
        "muted": RGBColor(0x70, 0x70, 0x90),
    },
    "light_clean": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x1A, 0x1A, 0x2E),
        "body": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x2D, 0x5B, 0xE3),
        "subtitle": RGBColor(0x55, 0x55, 0x77),
        "muted": RGBColor(0x99, 0x99, 0xAA),
    },
    "navy_gold": {
        "bg": RGBColor(0x0D, 0x1B, 0x2A),
        "title": RGBColor(0xF0, 0xE6, 0xD0),
        "body": RGBColor(0xC8, 0xC8, 0xD0),
        "accent": RGBColor(0xD4, 0xA5, 0x37),
        "subtitle": RGBColor(0xA0, 0x90, 0x70),
        "muted": RGBColor(0x60, 0x60, 0x80),
    },
    "minimal_gray": {
        "bg": RGBColor(0xF5, 0xF5, 0xF5),
        "title": RGBColor(0x22, 0x22, 0x22),
        "body": RGBColor(0x44, 0x44, 0x44),
        "accent": RGBColor(0xE0, 0x40, 0x40),
        "subtitle": RGBColor(0x66, 0x66, 0x66),
        "muted": RGBColor(0xAA, 0xAA, 0xAA),
    },
}

DEFAULT_THEME = "dark_academic"

# Theme detection keywords
_THEME_KEYWORDS = {
    "light_clean": ["밝은", "라이트", "light", "white", "clean", "흰색"],
    "navy_gold": ["네이비", "골드", "navy", "gold", "고급", "럭셔리"],
    "minimal_gray": ["미니멀", "그레이", "minimal", "gray", "심플"],
    "dark_academic": ["다크", "어두운", "dark", "학술"],
}


def detect_theme(instruction: str) -> str:
    """Detect theme from user instruction."""
    lower = instruction.lower()
    for theme_name, keywords in _THEME_KEYWORDS.items():
        if any(kw in lower for kw in keywords):
            return theme_name
    return DEFAULT_THEME


# ---------------------------------------------------------------------------
# Font detection
# ---------------------------------------------------------------------------
def _detect_font(text: str) -> str:
    """Detect appropriate font based on text content."""
    if re.search(r'[가-힣]', text):
        return "Noto Sans KR"
    return "Noto Sans"


def _set_font(paragraph, text: str, size: int, color: RGBColor,
              bold: bool = False, italic: bool = False):
    """Set text and font properties on a paragraph."""
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    paragraph.font.name = _detect_font(text)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def _set_slide_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs: Presentation, title: str, subtitle: str, theme: dict):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["bg"])

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_font(tf.paragraphs[0], title, 40, theme["title"], bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        _set_font(tf2.paragraphs[0], subtitle, 20, theme["subtitle"])
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def _add_section_slide(prs: Presentation, title: str, theme: dict):
    """Add a section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["bg"])

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_font(tf.paragraphs[0], title, 36, theme["accent"], bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def _add_content_slide(prs: Presentation, title: str, bullets: list[str],
                        notes: str, visual: str, theme: dict,
                        image_bytes: Optional[bytes] = None):
    """Add a content slide with optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["bg"])

    # Determine layout: with image → 2-column, without → full width
    has_image = image_bytes is not None
    content_width = Inches(6) if has_image else Inches(11.5)

    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_font(tf.paragraphs[0], title, 28, theme["title"], bold=True)

    # Accent line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(2.5), Emu(36000))
    line.fill.solid()
    line.fill.fore_color.rgb = theme["accent"]
    line.line.fill.background()

    # Content bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), content_width, Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        bullet_text = f"• {bullet}"
        _set_font(p, bullet_text, 18, theme["body"])
        p.space_after = Pt(8)

    # Image (right side)
    if has_image:
        img_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            img_stream, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.5)
        )
    elif visual:
        # Visual suggestion text (small, at bottom)
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf3 = txBox3.text_frame
        _set_font(tf3.paragraphs[0], f"📊 {visual}", 10, theme["muted"], italic=True)

    # Speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


# ---------------------------------------------------------------------------
# JSON parsing
# ---------------------------------------------------------------------------
def parse_slides_json(raw: str) -> Optional[dict]:
    """Parse LLM output into slide structure dict."""
    text = raw.strip()

    # Strip markdown code fences
    if text.startswith("```"):
        first_nl = text.index("\n")
        text = text[first_nl + 1:]
    if text.endswith("```"):
        text = text[:-3].rstrip()

    start = text.find("{")
    end = text.rfind("}") + 1
    if start == -1 or end == 0:
        logger.warning("No JSON object found in presentation response")
        return None

    try:
        return json.loads(text[start:end])
    except json.JSONDecodeError as e:
        logger.warning(f"Failed to parse presentation JSON: {e}")
        return None


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------
def build_pptx(data: dict, theme_name: Optional[str] = None,
               images: Optional[list[Optional[bytes]]] = None) -> bytes:
    """Build a .pptx file from parsed slide data.

    Args:
        data: Dict with 'title', 'subtitle', 'slides' keys.
        theme_name: Theme name (dark_academic, light_clean, navy_gold, minimal_gray).
        images: Optional list of PNG bytes per slide (None for no image).

    Returns:
        Bytes of the .pptx file.
    """
    theme = THEMES.get(theme_name or DEFAULT_THEME, THEMES[DEFAULT_THEME])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    _add_title_slide(prs, data.get("title", ""), data.get("subtitle", ""), theme)

    # Content slides
    slides_data = data.get("slides", [])
    for i, slide_data in enumerate(slides_data):
        layout = slide_data.get("layout", "title_and_content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", [])
        notes = slide_data.get("notes", "")
        visual = slide_data.get("visual", "")

        # Get image for this slide if available
        img = images[i] if images and i < len(images) else None

        if layout == "title_slide":
            _add_title_slide(prs, title, notes, theme)
        elif layout == "section_header":
            s = _add_section_slide(prs, title, theme)
            if notes:
                s.notes_slide.notes_text_frame.text = notes
        else:
            if isinstance(content, str):
                content = [content]
            _add_content_slide(prs, title, content, notes, visual, theme, img)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Markdown preview
# ---------------------------------------------------------------------------
def format_slides_markdown(data: dict, theme_name: str = DEFAULT_THEME) -> str:
    """Format slide data as markdown preview."""
    lines = [f"# {data.get('title', 'Untitled')}"]
    if data.get("subtitle"):
        lines.append(f"\n*{data['subtitle']}*")
    lines.append(f"\n🎨 테마: `{theme_name}`\n")

    for i, slide in enumerate(data.get("slides", []), 1):
        lines.append(f"---\n\n### Slide {i}: {slide.get('title', '')}\n")

        content = slide.get("content", [])
        if isinstance(content, str):
            # Prose mode: render as paragraphs
            lines.append(content)
        else:
            # Bullet mode: render as list items
            for bullet in content:
                lines.append(f"- {bullet}")

        if slide.get("notes"):
            note_text = slide["notes"]
            preview = note_text[:100] + "..." if len(note_text) > 100 else note_text
            lines.append(f"\n> 🎤 *{preview}*")

        if slide.get("visual"):
            lines.append(f"\n📊 {slide['visual']}")

        lines.append("")

    return "\n".join(lines)
